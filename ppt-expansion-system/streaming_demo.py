"""
PPT 智能扩展系统 - 完整使用示例
包含流式输出、向量数据库管理、配置管理等
"""

import os
import asyncio
from typing import List, AsyncIterator, Dict
from pathlib import Path
import json

from langchain_core.callbacks import AsyncCallbackHandler
from langchain_core.outputs import LLMResult
from langchain_community.document_loaders import PyPDFLoader, TextLoader
from langchain.text_splitter import RecursiveCharacterTextSplitter


# ==================== 流式输出处理 ====================
class StreamingCallbackHandler(AsyncCallbackHandler):
    """流式输出回调处理器"""
    
    def __init__(self):
        self.tokens = []
        self.is_streaming = False
    
    async def on_llm_new_token(self, token: str, **kwargs) -> None:
        """处理新 token"""
        self.tokens.append(token)
        print(token, end="", flush=True)
    
    async def on_llm_end(self, response: LLMResult, **kwargs) -> None:
        """LLM 结束时"""
        print("\n")
        self.is_streaming = False


# ==================== 向量数据库管理器 ====================
class KnowledgeBaseManager:
    """知识库管理器 - 管理本地 RAG 向量数据库"""
    
    def __init__(self, db_path: str = "./knowledge_base"):
        self.db_path = Path(db_path)
        self.db_path.mkdir(parents=True, exist_ok=True)
        self.metadata_file = self.db_path / "metadata.json"
        self.metadata = self._load_metadata()
    
    def _load_metadata(self) -> Dict:
        """加载元数据"""
        if self.metadata_file.exists():
            with open(self.metadata_file, 'r', encoding='utf-8') as f:
                return json.load(f)
        return {"documents": [], "last_updated": None}
    
    def _save_metadata(self):
        """保存元数据"""
        with open(self.metadata_file, 'w', encoding='utf-8') as f:
            json.dump(self.metadata, f, ensure_ascii=False, indent=2)
    
    def add_documents_from_folder(self, folder_path: str) -> List:
        """从文件夹添加文档"""
        from langchain_core.documents import Document
        
        documents = []
        folder = Path(folder_path)
        
        # 支持的文件类型
        for file_path in folder.rglob("*"):
            if file_path.suffix in ['.txt', '.md', '.pdf']:
                try:
                    if file_path.suffix == '.pdf':
                        loader = PyPDFLoader(str(file_path))
                    else:
                        loader = TextLoader(str(file_path), encoding='utf-8')
                    
                    docs = loader.load()
                    
                    # 分割文档
                    text_splitter = RecursiveCharacterTextSplitter(
                        chunk_size=1000,
                        chunk_overlap=200
                    )
                    split_docs = text_splitter.split_documents(docs)
                    
                    documents.extend(split_docs)
                    
                    # 记录元数据
                    self.metadata["documents"].append({
                        "file": str(file_path),
                        "chunks": len(split_docs)
                    })
                    
                except Exception as e:
                    print(f"加载文件 {file_path} 失败: {e}")
        
        self._save_metadata()
        return documents
    
    def add_ppt_history(self, ppt_id: str, final_notes: str):
        """添加历史 PPT 扩展结果"""
        from langchain_core.documents import Document
        
        doc = Document(
            page_content=final_notes,
            metadata={
                "ppt_id": ppt_id,
                "type": "ppt_expansion_history"
            }
        )
        
        self.metadata["documents"].append({
            "ppt_id": ppt_id,
            "type": "history"
        })
        self._save_metadata()
        
        return [doc]
    
    def get_stats(self) -> Dict:
        """获取统计信息"""
        return {
            "total_documents": len(self.metadata["documents"]),
            "last_updated": self.metadata.get("last_updated"),
            "db_size": sum(f.stat().st_size for f in self.db_path.rglob("*") if f.is_file())
        }


# ==================== 配置管理器 ====================
class ConfigManager:
    """配置管理器"""
    
    def __init__(self, config_file: str = "config.json"):
        self.config_file = Path(config_file)
        self.config = self._load_config()
    
    def _load_config(self) -> Dict:
        """加载配置"""
        if self.config_file.exists():
            with open(self.config_file, 'r') as f:
                return json.load(f)
        
        # 默认配置
        return {
            "llm": {
                "api_key": os.getenv("OPENAI_API_KEY", ""),
                "base_url": os.getenv("OPENAI_BASE_URL", "https://api.openai.com/v1"),
                "model": "gpt-4"
            },
            "retrieval": {
                "preferred_sources": ["arxiv", "wikipedia"],
                "max_results": 3
            },
            "expansion": {
                "max_revisions": 2,
                "min_gap_priority": 3
            },
            "streaming": {
                "enabled": True,
                "chunk_size": 50
            }
        }
    
    def save_config(self):
        """保存配置"""
        with open(self.config_file, 'w') as f:
            json.dump(self.config, f, indent=2)
    
    def get(self, key: str, default=None):
        """获取配置项"""
        keys = key.split('.')
        value = self.config
        for k in keys:
            value = value.get(k, default)
            if value is None:
                return default
        return value


# ==================== 完整流程示例 ====================
class PPTExpansionPipeline:
    """PPT 扩展流水线 - 完整流程封装"""
    
    def __init__(self, config_file: str = "config.json"):
        self.config_manager = ConfigManager(config_file)
        self.kb_manager = KnowledgeBaseManager()
        
        # 初始化 LLM 配置
        from ppt_expansion_system import LLMConfig, PPTExpansionGraph
        self.llm_config = LLMConfig(
            api_key=self.config_manager.get("llm.api_key"),
            base_url=self.config_manager.get("llm.base_url"),
            model=self.config_manager.get("llm.model")
        )
        
        # 创建主 Graph
        self.graph = PPTExpansionGraph(self.llm_config)
        
        # 初始化向量数据库
        self._initialize_vectorstore()
    
    def _initialize_vectorstore(self):
        """初始化向量数据库"""
        # 从知识库文件夹加载文档
        docs = self.kb_manager.add_documents_from_folder("./knowledge_sources")
        if docs:
            self.graph.retrieval_agent.initialize_vectorstore(docs)
            print(f"✓ 向量数据库初始化完成，加载 {len(docs)} 个文档片段")
    
    def load_ppt(self, file_path: str) -> List[str]:
        """加载 PPT 文件"""
        # 这里需要集成 PPT 解析库，如 python-pptx
        # 示例：返回每一页的文本
        try:
            from pptx import Presentation
            prs = Presentation(file_path)
            
            ppt_texts = []
            for i, slide in enumerate(prs.slides):
                text_content = []
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text_content.append(shape.text)
                ppt_texts.append(f"Page {i+1}:\n" + "\n".join(text_content))
            
            return ppt_texts
        except ImportError:
            print("请安装 python-pptx: pip install python-pptx")
            return []
        except Exception as e:
            print(f"加载 PPT 失败: {e}")
            return []
    
    def run(self, ppt_texts: List[str], streaming: bool = True) -> Dict:
        """运行完整流程"""
        print("=" * 60)
        print("PPT 智能扩展系统 - 开始处理")
        print("=" * 60)
        
        # Step 0: 全局分析
        print("\n[Step 0] 全局结构分析与知识点划分...")
        
        # 运行主流程
        max_revisions = self.config_manager.get("expansion.max_revisions", 2)
        result = self.graph.run(ppt_texts, max_revisions=max_revisions)
        
        # 保存历史
        ppt_id = f"ppt_{len(self.kb_manager.metadata['documents']) + 1}"
        self.kb_manager.add_ppt_history(ppt_id, result["final_notes"])
        
        # 返回结果
        return {
            "ppt_id": ppt_id,
            "global_outline": result["global_outline"],
            "knowledge_units": result["knowledge_units"],
            "final_notes": result["final_notes"],
            "revision_count": result.get("revision_count", 0),
            "stats": {
                "total_gaps": len(result["knowledge_gaps"]),
                "expanded_items": len(result["expanded_content"]),
                "retrieved_docs": len(result["retrieved_docs"])
            }
        }
    
    async def run_streaming(self, ppt_texts: List[str]) -> AsyncIterator[str]:
        """流式运行（异步生成器）"""
        # TODO: 实现真正的流式输出
        # 这需要 LangChain 的异步流式支持
        result = self.run(ppt_texts, streaming=False)
        
        # 模拟流式输出
        content = result["final_notes"]
        chunk_size = self.config_manager.get("streaming.chunk_size", 50)
        
        for i in range(0, len(content), chunk_size):
            chunk = content[i:i+chunk_size]
            yield chunk
            await asyncio.sleep(0.05)  # 模拟延迟
    
    def export_to_markdown(self, result: Dict, output_file: str):
        """导出为 Markdown 文件"""
        output_path = Path(output_file)
        
        with open(output_path, 'w', encoding='utf-8') as f:
            f.write(f"# PPT 扩展笔记 - {result['ppt_id']}\n\n")
            f.write(f"## 全局结构\n\n")
            f.write(f"```json\n{json.dumps(result['global_outline'], ensure_ascii=False, indent=2)}\n```\n\n")
            f.write(f"## 知识单元\n\n")
            for unit in result["knowledge_units"]:
                f.write(f"### {unit.title}\n")
                f.write(f"- 涉及页面: {', '.join(map(str, unit.pages))}\n")
                f.write(f"- 核心概念: {', '.join(unit.core_concepts)}\n\n")
            f.write(f"## 扩展笔记\n\n")
            f.write(result["final_notes"])
            f.write(f"\n\n---\n")
            f.write(f"生成统计: {result['stats']}\n")
        
        print(f"✓ 结果已导出到: {output_path}")


# ==================== 主程序 ====================
async def main():
    """主程序"""
    # 1. 初始化流水线
    pipeline = PPTExpansionPipeline("config.json")
    
    # 2. 准备 PPT 文本（可以从文件加载）
    ppt_texts = [
        """第1页：Transformer 架构概述
        - 基于注意力机制的神经网络架构
        - 2017年由Google提出
        - 应用于机器翻译、文本生成等任务""",
        
        """第2页：Self-Attention 机制
        公式：Attention(Q, K, V) = softmax(QK^T / √d_k)V
        其中：
        - Q = XW_q（查询矩阵）
        - K = XW_k（键矩阵）
        - V = XW_v（值矩阵）""",
        
        """第3页：Position Encoding
        为什么需要：Transformer没有循环结构，无法捕获位置信息
        实现方式：
        PE(pos, 2i) = sin(pos / 10000^(2i/d_model))
        PE(pos, 2i+1) = cos(pos / 10000^(2i/d_model))""",
        
        """第4页：Multi-Head Attention
        - 使用多个注意力头并行计算
        - 每个头关注不同的表示子空间
        - 最后拼接所有头的输出"""
    ]
    
    # 3. 运行流程（非流式）
    print("\n选择运行模式：")
    print("1. 标准模式（完整处理）")
    print("2. 流式模式（实时输出）")
    
    mode = input("请选择 (1/2): ").strip()
    
    if mode == "2":
        # 流式模式
        print("\n开始流式处理...\n")
        async for chunk in pipeline.run_streaming(ppt_texts):
            print(chunk, end="", flush=True)
        print("\n\n流式处理完成！")
    else:
        # 标准模式
        result = pipeline.run(ppt_texts, streaming=False)
        
        # 4. 打印结果
        print("\n" + "=" * 60)
        print("处理完成！")
        print("=" * 60)
        print(f"\nPPT ID: {result['ppt_id']}")
        print(f"修订次数: {result['revision_count']}")
        print(f"\n统计信息:")
        print(f"  - 识别知识缺口: {result['stats']['total_gaps']} 个")
        print(f"  - 生成扩展内容: {result['stats']['expanded_items']} 项")
        print(f"  - 检索参考文档: {result['stats']['retrieved_docs']} 篇")
        
        print(f"\n最终笔记预览:")
        print("-" * 60)
        print(result["final_notes"][:500] + "..." if len(result["final_notes"]) > 500 else result["final_notes"])
        
        # 5. 导出结果
        export = input("\n是否导出为 Markdown? (y/n): ").strip().lower()
        if export == 'y':
            output_file = f"output_{result['ppt_id']}.md"
            pipeline.export_to_markdown(result, output_file)
        
        # 6. 显示知识库统计
        kb_stats = pipeline.kb_manager.get_stats()
        print(f"\n知识库统计:")
        print(f"  - 总文档数: {kb_stats['total_documents']}")
        print(f"  - 数据库大小: {kb_stats['db_size'] / 1024:.2f} KB")


# ==================== 命令行工具 ====================
if __name__ == "__main__":
    import sys
    
    if len(sys.argv) > 1:
        # 命令行模式
        ppt_file = sys.argv[1]
        pipeline = PPTExpansionPipeline()
        
        ppt_texts = pipeline.load_ppt(ppt_file)
        if ppt_texts:
            result = pipeline.run(ppt_texts)
            pipeline.export_to_markdown(result, f"output_{Path(ppt_file).stem}.md")
    else:
        # 交互模式
        asyncio.run(main())