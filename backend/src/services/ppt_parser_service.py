from typing import Any, Dict, List, Optional
import base64
import os
import shutil
import subprocess
import tempfile
import fitz  # PyMuPDF
from pptx import Presentation
from pptx.enum.shapes import MSO_SHAPE_TYPE


class DocumentParserService:
    """Parse PPTX/PDF into structured slide items with semantic hierarchy."""

    def __init__(self):
        pass

    def parse_document(self, path: str, ext: str) -> List[Dict[str, Any]]:
        if ext == ".pdf":
            return self._parse_pdf(path)
        if ext == ".pptx":
            pdf_path = self._convert_pptx_to_pdf(path)
            try:
                return self._parse_pdf(pdf_path)
            finally:
                if pdf_path and os.path.exists(pdf_path):
                    try:
                        os.remove(pdf_path)
                    except Exception as e:
                        print(f"⚠️ 清理临时PDF文件失败: {e}")
        return []

    def _convert_pptx_to_pdf(self, pptx_path: str) -> str:
        """将 PPTX 文件转换为 PDF"""
        libreoffice_paths = [
            "soffice",  
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",  
            "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe",  
        ]
        
        soffice = None
        for path in libreoffice_paths:
            if os.path.exists(path) or path == "soffice":
                try:
                    result = subprocess.run(
                        [path, "--version"],
                        capture_output=True,
                        timeout=5
                    )
                    if result.returncode == 0:
                        soffice = path
                        break
                except:
                    continue
        
        if not soffice:
            raise Exception("未找到LibreOffice，无法将PPTX转换为PDF。请安装LibreOffice：https://www.libreoffice.org/")
        
        temp_dir = tempfile.mkdtemp()
        try:
            # 使用LibreOffice将PPTX转换为PDF
            abs_path = os.path.abspath(pptx_path)
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, abs_path],
                capture_output=True,
                timeout=60000  
            )
            
            if result.returncode != 0:
                error_msg = result.stderr.decode() if result.stderr else "未知错误"
                raise Exception(f"LibreOffice转换失败: {error_msg}")
            
            pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
            pdf_path = os.path.join(temp_dir, pdf_name)
            
            if not os.path.exists(pdf_path):
                raise Exception("LibreOffice未生成PDF文件")
            
            final_pdf_path = tempfile.NamedTemporaryFile(suffix='.pdf', delete=False)
            final_pdf_path.close()
            
            shutil.copy2(pdf_path, final_pdf_path.name)

            shutil.rmtree(temp_dir, ignore_errors=True)
            
            return final_pdf_path.name
        except subprocess.TimeoutExpired:
            raise Exception("PPTX转PDF超时，文件可能过大或LibreOffice无响应")
        except Exception as e:
            if 'temp_dir' in locals():
                shutil.rmtree(temp_dir, ignore_errors=True)
            raise

    def _render_page_to_image(self, page, zoom: float = 2.0) -> str:
        """将PDF页面渲染为base64编码的PNG图片"""
        mat = fitz.Matrix(zoom, zoom)
        pix = page.get_pixmap(matrix=mat)
        img_data = pix.tobytes("png")
        base64_str = base64.b64encode(img_data).decode('utf-8')
        return f"data:image/png;base64,{base64_str}"

    def _render_pptx_slide_to_image(self, pptx_path: str, slide_index: int) -> Optional[str]:
        """将PPTX的某一页渲染为base64编码的PNG图片"""
        try:
            # 方法1: 尝试使用comtypes调用PowerPoint (Windows)
            if os.name == 'nt':  
                try:
                    return self._render_pptx_with_powerpoint(pptx_path, slide_index)
                except Exception:
                    pass
            
            # 方法2: 尝试使用LibreOffice
            try:
                return self._render_pptx_with_libreoffice(pptx_path, slide_index)
            except Exception:
                pass
            
            # 如果都失败，返回None，前端会显示占位符
            return None
        except Exception as e:
            print(f"渲染PPTX幻灯片失败: {e}")
            return None

    def _render_pptx_with_powerpoint(self, pptx_path: str, slide_index: int) -> str:
        """使用PowerPoint COM接口渲染PPTX (仅Windows)"""
        try:
            import comtypes.client
            import comtypes.gen.Microsoft.Office.Interop.PowerPoint as PPT
            
            powerpoint = comtypes.client.CreateObject("PowerPoint.Application")
            powerpoint.Visible = 0 
            
            abs_path = os.path.abspath(pptx_path)
            presentation = powerpoint.Presentations.Open(abs_path)
            
            # 获取指定幻灯片
            slide = presentation.Slides[slide_index]
            
            # 导出为图片
            temp_img = tempfile.NamedTemporaryFile(suffix='.png', delete=False)
            temp_img_path = temp_img.name
            temp_img.close()
            
            # 导出幻灯片为PNG 
            slide.Export(temp_img_path, "PNG", 1920, 1080)
            
            # 读取图片并转换为base64
            with open(temp_img_path, 'rb') as f:
                img_data = f.read()
            base64_str = base64.b64encode(img_data).decode('utf-8')
            
            presentation.Close()
            powerpoint.Quit()
            os.unlink(temp_img_path)
            
            return f"data:image/png;base64,{base64_str}"
        except ImportError:
            raise Exception("需要安装comtypes: pip install comtypes")
        except Exception as e:
            raise Exception(f"PowerPoint渲染失败: {e}")

    def _render_pptx_with_libreoffice(self, pptx_path: str, slide_index: int) -> str:
        """使用LibreOffice命令行工具渲染PPTX"""
        libreoffice_paths = [
            "soffice",  
            "C:\\Program Files\\LibreOffice\\program\\soffice.exe",  
            "C:\\Program Files (x86)\\LibreOffice\\program\\soffice.exe", 
        ]
        
        soffice = None
        for path in libreoffice_paths:
            if os.path.exists(path) or path == "soffice":
                try:
                    result = subprocess.run(
                        [path, "--version"],
                        capture_output=True,
                        timeout=5
                    )
                    if result.returncode == 0:
                        soffice = path
                        break
                except:
                    continue
        
        if not soffice:
            raise Exception("未找到LibreOffice，请安装LibreOffice")
        
        with tempfile.TemporaryDirectory() as temp_dir:
            # 使用LibreOffice将PPTX转换为PDF
            abs_path = os.path.abspath(pptx_path)
            result = subprocess.run(
                [soffice, "--headless", "--convert-to", "pdf", "--outdir", temp_dir, abs_path],
                capture_output=True,
                timeout=30
            )
            
            if result.returncode != 0:
                raise Exception(f"LibreOffice转换失败: {result.stderr.decode()}")
            
            pdf_name = os.path.splitext(os.path.basename(pptx_path))[0] + ".pdf"
            pdf_path = os.path.join(temp_dir, pdf_name)
            
            if not os.path.exists(pdf_path):
                raise Exception("LibreOffice未生成PDF文件")
            
            # 使用PyMuPDF渲染PDF的指定页面
            pdf_doc = fitz.open(pdf_path)
            if slide_index < len(pdf_doc):
                page = pdf_doc[slide_index]
                img_base64 = self._render_page_to_image(page)
                pdf_doc.close()
                return img_base64
            else:
                pdf_doc.close()
                raise Exception(f"幻灯片索引 {slide_index} 超出范围")

    def _parse_pdf(self, path: str) -> List[Dict[str, Any]]:
        doc = fitz.open(path)
        slides: List[Dict[str, Any]] = []

        for i, page in enumerate(doc, start=1):
            page_dict = page.get_text("dict")
            blocks = page_dict["blocks"]

            text_items = []
            image_count = 0

            for b in blocks:
                if b["type"] == 0:  # Text block
                    for line in b["lines"]:
                        for span in line["spans"]:
                            txt = span["text"].strip()
                            # 过滤掉微小噪点
                            if txt and len(txt) > 0 and span["size"] > 5:
                                text_items.append({
                                    "text": txt,
                                    "size": span["size"],
                                    "bbox": span["bbox"]
                                })
                elif b["type"] == 1:  # Image block
                    image_count += 1
            # 按垂直位置排序
            text_items.sort(key=lambda x: x["bbox"][1])

            title = ""
            raw_points = []
            slide_type = "content"

            if text_items:
                max_size = max(t["size"] for t in text_items)
                
                # 1. 提取标题 
                # 第一优先级：字号最大且在页面上方的文字
                title_candidates = [
                    t for t in text_items
                    if t["size"] >= max_size * 0.9 and t["bbox"][1] < page.rect.height * 0.4
                ]
                
                # 第二优先级：字号最大
                if not title_candidates:
                    title_candidates = [
                        t for t in text_items
                        if t["size"] >= max_size * 0.85 and t["bbox"][1] < page.rect.height * 0.5
                    ]
                
                # 第三优先级：最大字号的第一个
                if not title_candidates and text_items:
                    title_candidates = [t for t in text_items if t["size"] == max_size][:1]
                
                # 合并标题候选
                if title_candidates:
                    title = " ".join([t["text"].strip() for t in title_candidates if t["text"].strip()])
                    # 过滤掉纯数字
                    if title and not title.replace(" ", "").isdigit():
                        pass  
                    else:
                        title = ""  
                
                if not title:
                    meaningful_items = [
                        t for t in text_items
                        if t["text"].strip() 
                        and len(t["text"].strip()) >= 2  
                        and not t["text"].strip().replace(" ", "").isdigit()  
                        and t["size"] > 10 
                    ]
                    
                    if meaningful_items:
                        meaningful_items.sort(key=lambda x: (-x["size"], x["bbox"][1]))
                        title = meaningful_items[0]["text"].strip()
                
                if not title:
                    title = f"Page {i}"

                # 2. 提取并合并正文内容
                body_items = [
                    t for t in text_items
                    if t not in title_candidates and t["size"] > 8 
                ]

                merged_points = []
                if body_items:
                    current_point = body_items[0]["text"]
                    current_size = body_items[0]["size"]
                    current_bbox = body_items[0]["bbox"]

                    for item in body_items[1:]:
                        font_diff = abs(item["size"] - current_size)
                        vertical_dist = item["bbox"][1] - current_bbox[3] 
                        is_list_start = item["text"].strip().startswith(('•', '-', '▪', '➢', '1.', '2.', '(', '（'))

                        line_height_threshold = item["size"] * 1.5 
                        
                        if font_diff < 1.0 and vertical_dist < line_height_threshold and not is_list_start:
                             current_point += ("" if item["text"].startswith('，') or item["text"].startswith('。') else "") + item["text"]
                             current_bbox = item["bbox"] 
                        else:
                            if len(current_point) > 2: 
                                merged_points.append(current_point)
                            current_point = item["text"]
                            current_size = item["size"]
                            current_bbox = item["bbox"]
                    
                    if len(current_point) > 2:
                        merged_points.append(current_point)
                
                raw_points = merged_points

            # 语义类型推断
            lower_title = title.lower()
            if "overview" in lower_title or "contents" in lower_title or "目录" in lower_title:
                slide_type = "toc"
            elif image_count > 0 and len(raw_points) < 3:
                slide_type = "figure"

            if title.startswith("Page "):
                for point in raw_points:
                    if point.strip() and len(point.strip()) >= 2 and not point.strip().isdigit():
                        title = point.strip()[:100]  
                        break
            
            # 统一数据结构
            structured_points = [
                {"type": "text", "level": 0, "text": p} for p in raw_points
            ]

            # 渲染页面为图片
            page_image = self._render_page_to_image(page)
            
            slides.append({
                "page_num": i,
                "title": title[:100],
                "type": slide_type,
                "raw_points": structured_points[:15],
                "images": ([f"[包含 {image_count} 张图片/图表]"] if image_count > 0 else []),
                "expanded_html": "<p><i>待补充 AI 深度解析内容...</i></p>",
                "references": [],
                "image": page_image 
            })
        doc.close()
        return slides

    def _parse_pptx(self, path: str) -> List[Dict[str, Any]]:
        prs = Presentation(path)
        slides: List[Dict[str, Any]] = []
        
        for i, slide in enumerate(prs.slides, start=1):
            # 1. 布局语义分析
            layout_name = slide.slide_layout.name
            slide_type = "content"
            if "title slide" in layout_name.lower() or i == 1:
                slide_type = "cover"
            elif "section" in layout_name.lower():
                slide_type = "section"
            
            title = ""
            subtitle = ""
            bullets = []
            img_descriptions = []
            
            # 2. 遍历形状提取内容
            sorted_shapes = sorted(
                slide.shapes, 
                key=lambda x: (x.top if hasattr(x, 'top') else 0, x.left if hasattr(x, 'left') else 0)
            )

            for shape in sorted_shapes:
                # 提取图片描述 
                if shape.shape_type == MSO_SHAPE_TYPE.PICTURE:
                    if hasattr(shape, "image"):
                         desc = f"Image: {shape.name}"
                         try:
                             desc += f" ({shape.width//9525}x{shape.height//9525} px)"
                         except:
                             pass
                         img_descriptions.append(desc)

                # 提取文本内容 
                if shape.has_text_frame:
                    if shape == slide.shapes.title:
                        title = shape.text
                        continue
                        
                    for paragraph in shape.text_frame.paragraphs:
                        txt = paragraph.text.strip()
                        if txt:
                            bullets.append({
                                "type": "text",
                                "level": paragraph.level,
                                "text": txt
                            })

                # 提取表格内容
                if shape.has_table:
                    table_data = []
                    for row in shape.table.rows:
                        row_data = []
                        for cell in row.cells:
                            cell_text = cell.text_frame.text.strip()
                            row_data.append(cell_text)
                        table_data.append(row_data)
                    
                    if any(any(cell for cell in row) for row in table_data): 
                        bullets.append({
                            "type": "table",
                            "data": table_data
                        })

            # 3. 语义推断补充
            if not title and bullets and bullets[0]["type"] == "text":
                title = bullets[0]["text"]

            if slide_type == "content" and len(bullets) > 8:
                 pass

            # 渲染幻灯片为图片 
            slide_image = self._render_pptx_slide_to_image(path, i - 1)
            
            slides.append({
                "page_num": i,
                "title": title if title else f"Slide {i}",
                "type": slide_type,
                "raw_points": bullets[:20], 
                "images": img_descriptions,
                "expanded_html": "<p><i>待补充 AI 深度解析内容...</i></p>",
                "references": [],
                "image": slide_image  
            })
        
        return slides
